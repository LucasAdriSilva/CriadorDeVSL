import collections , collections.abc, sys, os, base64, pyautogui
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN ,MSO_ANCHOR
from datetime import datetime
from time import sleep
import GerarCodigo as Gc

Gc.gerartudo()

def resourcePath(path): 
    try:
        basicPath = sys.__MEIPASS
    except Exception:
        basicPath = os.path.abspath('.')
    return os.path.join(basicPath, path)

try:
    #Abrindo PowerPoint
    apresentacao = Presentation(resourcePath('pptx\\template\\default.pptx'))


    #Lendo o arquivo.tx
    with open("copy.txt", "r", encoding='UTF-8') as tf:
        lines = tf.read().split('\n')
        
        # Cada linah do arquivo.txt ira fazer esse bloco
        for line in lines:

            #Criando um slide
            slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])

            #Criando e centralizando a TextBox
            x = Inches(1)
            y = Inches(3)
            largura = Inches(8)
            altura = Inches(1)
            caixa_texto = slide.shapes.add_textbox(x, y, largura, altura)

            #Formatando o TextBox
            caixa_texto.text = line
            caixa_texto.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            caixa_texto.text_frame.paragraphs[0].font.size = Pt(40)
            caixa_texto.text_frame.paragraphs[0].font.name = "Arial"
            caixa_texto.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


    #Salvando o arquivo pptx
    apresentacao.save("VSL.pptx")
    tf.close()
    pyautogui.alert(text=f'Seu arquivo foi criado com sucesso às {datetime.now().hour}:{datetime.now().minute}', title='Confirmação', button='Ok')

except Exception as e:
    print(e)
    sleep(20)